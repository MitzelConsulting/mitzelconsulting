#!/usr/bin/env python3
"""
Ingest only PPTX files (PowerPoint presentations).
"""

import os
import logging
from dotenv import load_dotenv
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from openai import OpenAI
from pinecone import Pinecone
from pptx import Presentation
import io

load_dotenv('.env.local')

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('ingest_pptx_only.log')
    ]
)
logger = logging.getLogger(__name__)

def setup_clients():
    """Initialize clients."""
    logger.info("🔧 Setting up clients...")
    
    credentials = service_account.Credentials.from_service_account_file(
        'service-account.json',
        scopes=['https://www.googleapis.com/auth/drive.readonly']
    )
    drive_service = build('drive', 'v3', credentials=credentials)
    
    openai_client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
    
    pc = Pinecone(api_key=os.getenv('PINECONE_API_KEY'))
    index = pc.Index('mizelconsulting')
    
    logger.info("✅ Clients initialized")
    return drive_service, openai_client, index

def get_pptx_files(drive_service, folder_id):
    """Get only PPTX files from the folder."""
    logger.info(f"📁 Getting PPTX files from folder: {folder_id}")
    
    try:
        results = drive_service.files().list(
            q=f"'{folder_id}' in parents and trashed=false and mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'",
            fields="files(id, name, mimeType, modifiedTime, size)"
        ).execute()
        
        files = results.get('files', [])
        logger.info(f"📊 Found {len(files)} PPTX files")
        
        return files
        
    except Exception as e:
        logger.error(f"❌ Error getting files: {e}")
        return []

def download_file(drive_service, file_id):
    """Download file content from Google Drive."""
    try:
        request = drive_service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while done is False:
            status, done = downloader.next_chunk()
        
        return fh.getvalue()
    except Exception as e:
        logger.error(f"❌ Error downloading file: {e}")
        return None

def extract_text_from_pptx(content):
    """Extract text from PPTX content."""
    try:
        prs = Presentation(io.BytesIO(content))
        
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + " "
            
            if slide_text.strip():
                text += f"Slide {slide_num}: {slide_text.strip()}\n\n"
        
        return text.strip()
    except Exception as e:
        logger.error(f"❌ Error extracting PPTX text: {e}")
        return ""

def process_pptx_file(file_metadata, drive_service, openai_client, index):
    """Process a single PPTX file and create one chunk."""
    file_id = file_metadata['id']
    file_name = file_metadata['name']
    
    logger.info(f"🔄 Processing: {file_name}")
    
    try:
        # Download file
        content = download_file(drive_service, file_id)
        if not content:
            logger.error(f"❌ Failed to download {file_name}")
            return False
        
        logger.info(f"✅ Downloaded {len(content)} bytes")
        
        # Extract text
        text = extract_text_from_pptx(content)
        
        if not text:
            logger.warning(f"⚠️ No text extracted from {file_name}")
            return False
        
        logger.info(f"✅ Extracted {len(text)} characters")
        
        # Take first 800 characters as one chunk
        chunk_text = text[:800] + "..." if len(text) > 800 else text
        logger.info(f"📄 Created chunk with {len(chunk_text)} characters")
        
        # Generate embedding
        logger.info("🧠 Generating embedding...")
        response = openai_client.embeddings.create(
            model="text-embedding-3-small",
            input=chunk_text
        )
        
        embedding = response.data[0].embedding
        logger.info(f"✅ Generated embedding with {len(embedding)} dimensions")
        
        # Prepare vector
        vector = {
            'id': f"{file_id}_pptx",
            'values': embedding,
            'metadata': {
                'text': chunk_text,
                'file_id': file_id,
                'file_name': file_name,
                'folder_path': 'site',
                'mime_type': file_metadata.get('mimeType', ''),
                'modified_time': file_metadata.get('modifiedTime', '')
            }
        }
        
        # Upsert to Pinecone
        logger.info("📤 Upserting to Pinecone...")
        index.upsert(vectors=[vector], namespace='site')
        logger.info("✅ Successfully upserted to Pinecone")
        
        return True
        
    except Exception as e:
        logger.error(f"❌ Error processing {file_name}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return False

def test_query(index, openai_client):
    """Test query to verify ingestion worked."""
    logger.info("🔍 Testing query...")
    
    try:
        query = "safety training presentation"
        response = openai_client.embeddings.create(
            model="text-embedding-3-small",
            input=query
        )
        query_embedding = response.data[0].embedding
        
        results = index.query(
            vector=query_embedding,
            top_k=5,
            namespace='site',
            include_metadata=True
        )
        
        logger.info(f"🔍 Query returned {len(results.matches)} results:")
        for match in results.matches:
            logger.info(f"  - {match.metadata['file_name']} (Score: {match.score:.4f})")
        
        return True
        
    except Exception as e:
        logger.error(f"❌ Error testing query: {e}")
        return False

def main():
    """Main function."""
    logger.info("🚀 Starting PPTX-only ingestion")
    
    try:
        drive_service, openai_client, index = setup_clients()
        
        folder_id = os.getenv('GDRIVE_FOLDER_ID')
        files = get_pptx_files(drive_service, folder_id)
        
        if not files:
            logger.error("❌ No PPTX files found")
            return
        
        successful_uploads = 0
        failed_uploads = 0
        
        for i, file_metadata in enumerate(files, 1):
            logger.info(f"\n{'='*60}")
            logger.info(f"📄 Processing file {i}/{len(files)}: {file_metadata['name']}")
            logger.info(f"{'='*60}")
            
            if process_pptx_file(file_metadata, drive_service, openai_client, index):
                successful_uploads += 1
                logger.info(f"✅ Successfully processed {file_metadata['name']}")
            else:
                failed_uploads += 1
                logger.error(f"❌ Failed to process {file_metadata['name']}")
        
        # Test query
        logger.info(f"\n{'='*60}")
        logger.info("🧪 Testing query functionality")
        logger.info(f"{'='*60}")
        test_query(index, openai_client)
        
        # Summary
        logger.info(f"\n{'='*60}")
        logger.info("📊 PPTX INGESTION SUMMARY")
        logger.info(f"{'='*60}")
        logger.info(f"📁 Files attempted: {len(files)}")
        logger.info(f"✅ Files successful: {successful_uploads}")
        logger.info(f"❌ Files failed: {failed_uploads}")
        logger.info(f"{'='*60}")
        
    except Exception as e:
        logger.error(f"❌ Pipeline failed: {e}")
        import traceback
        logger.error(traceback.format_exc())

if __name__ == "__main__":
    main()

