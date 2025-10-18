#!/usr/bin/env python3
"""
Comprehensive ingestion script for all supported file types:
- DOCX files (already working)
- PDF files
- PPT/PPTX files  
- DOC files
"""

import os
import logging
import subprocess
import tempfile
import re
from dotenv import load_dotenv
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from openai import OpenAI
from pinecone import Pinecone
from docx import Document
import io
import PyPDF2
from pptx import Presentation

load_dotenv('.env.local')

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('ingest_all_files.log')
    ]
)
logger = logging.getLogger(__name__)

def setup_clients():
    """Initialize clients."""
    logger.info("🔧 Setting up clients...")
    
    credentials = service_account.Credentials.from_service_account_file(
        'service-account.json',
        scopes=['https://www.googleapis.com/auth/drive.readonly']
    )
    drive_service = build('drive', 'v3', credentials=credentials)
    
    openai_client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
    
    pc = Pinecone(api_key=os.getenv('PINECONE_API_KEY'))
    index = pc.Index('mizelconsulting')
    
    logger.info("✅ Clients initialized")
    return drive_service, openai_client, index

def get_all_supported_files(drive_service, folder_id):
    """Get all supported files from the folder."""
    logger.info(f"📁 Getting all supported files from folder: {folder_id}")
    
    try:
        # Query for all supported file types
        supported_types = [
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # DOCX
            'application/pdf',  # PDF
            'application/vnd.ms-powerpoint',  # PPT
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # PPTX
            'application/msword'  # DOC
        ]
        
        # Build query for supported file types
        type_conditions = [f"mimeType='{mime_type}'" for mime_type in supported_types]
        type_query = " or ".join(type_conditions)
        full_query = f"'{folder_id}' in parents and trashed=false and ({type_query})"
        
        results = drive_service.files().list(
            q=full_query,
            fields="files(id, name, mimeType, modifiedTime, size)"
        ).execute()
        
        files = results.get('files', [])
        logger.info(f"📊 Found {len(files)} supported files")
        
        # Categorize files
        docx_files = []
        pdf_files = []
        ppt_files = []
        doc_files = []
        
        for file in files:
            mime_type = file.get('mimeType', '')
            if 'wordprocessingml' in mime_type:
                docx_files.append(file)
            elif 'pdf' in mime_type:
                pdf_files.append(file)
            elif 'presentation' in mime_type or 'powerpoint' in mime_type:
                ppt_files.append(file)
            elif 'msword' in mime_type:
                doc_files.append(file)
        
        logger.info(f"📄 DOCX files: {len(docx_files)}")
        logger.info(f"📄 PDF files: {len(pdf_files)}")
        logger.info(f"📊 PPT/PPTX files: {len(ppt_files)}")
        logger.info(f"📄 DOC files: {len(doc_files)}")
        
        return files
        
    except Exception as e:
        logger.error(f"❌ Error getting files: {e}")
        return []

def download_file(drive_service, file_id):
    """Download file content from Google Drive."""
    try:
        request = drive_service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while done is False:
            status, done = downloader.next_chunk()
        
        return fh.getvalue()
    except Exception as e:
        logger.error(f"❌ Error downloading file: {e}")
        return None

def extract_text_from_docx(content):
    """Extract text from DOCX content."""
    try:
        doc = Document(io.BytesIO(content))
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        return text.strip()
    except Exception as e:
        logger.error(f"❌ Error extracting DOCX text: {e}")
        return ""

def extract_text_from_pdf(content):
    """Extract text from PDF content."""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"❌ Error extracting PDF text: {e}")
        return ""

def extract_text_from_ppt(content, mime_type):
    """Extract text from PowerPoint files."""
    try:
        if 'presentationml' in mime_type:  # PPTX files
            prs = Presentation(io.BytesIO(content))
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += f"Slide {slide_num}: {shape.text.strip()}\n"
            return text.strip()
        else:  # PPT files (older format)
            # Try LibreOffice first, then fall back to binary extraction
            text = extract_text_from_old_ppt(content)
            if not text:
                text = extract_text_from_ppt_binary(content)
            return text
        
    except Exception as e:
        logger.error(f"❌ Error extracting PPT text: {e}")
        return ""

def extract_text_from_old_ppt(content):
    """Extract text from old PPT files using LibreOffice command line."""
    try:
        # Create a temporary file for the PPT content
        with tempfile.NamedTemporaryFile(suffix='.ppt', delete=False) as temp_file:
            temp_file.write(content)
            temp_file.flush()
            temp_ppt_path = temp_file.name
        
        try:
            # Use LibreOffice to convert PPT to text
            # LibreOffice command: soffice --headless --convert-to txt
            result = subprocess.run([
                '/Applications/LibreOffice.app/Contents/MacOS/soffice', '--headless', '--convert-to', 'txt', 
                temp_ppt_path
            ], capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                # LibreOffice creates the converted file in the current working directory
                # Extract just the filename from the temp path
                temp_filename = os.path.basename(temp_ppt_path)
                txt_filename = temp_filename.replace('.ppt', '.txt')
                txt_path = os.path.join(os.getcwd(), txt_filename)
                
                if os.path.exists(txt_path):
                    with open(txt_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                    # Clean up the text file
                    os.unlink(txt_path)
                    return text.strip()
                else:
                    logger.warning("⚠️ LibreOffice conversion succeeded but no text file found")
                    return ""
            else:
                logger.warning(f"⚠️ LibreOffice conversion failed: {result.stderr}")
                return ""
                
        finally:
            # Clean up the temporary PPT file
            if os.path.exists(temp_ppt_path):
                os.unlink(temp_ppt_path)
                
    except subprocess.TimeoutExpired:
        logger.warning("⚠️ LibreOffice conversion timed out")
        return ""
    except Exception as e:
        logger.warning(f"⚠️ Error with LibreOffice conversion: {e}")
        return ""

def extract_text_from_ppt_binary(content):
    """Extract readable text directly from PPT binary content."""
    try:
        logger.info("🧪 Trying binary text extraction for PPT...")
        
        # Convert to string and look for readable text patterns
        text_content = content.decode('latin-1', errors='ignore')
        
        # Look for common text patterns in PPT files
        patterns = [
            r'[\x20-\x7E]{15,}',  # Printable ASCII characters, at least 15 chars
            r'[A-Za-z\s]{25,}',   # Letters and spaces, at least 25 chars
        ]
        
        extracted_text = []
        for pattern in patterns:
            matches = re.findall(pattern, text_content)
            for match in matches:
                # Clean up the match
                clean_match = match.strip()
                # Filter out binary noise and keep meaningful text
                if (len(clean_match) > 15 and 
                    not any(char in clean_match for char in ['\x00', '\x01', '\x02', '\x03']) and
                    any(c.isalpha() for c in clean_match) and  # Must contain letters
                    len([c for c in clean_match if c.isalpha()]) > 5):  # At least 5 letters
                    extracted_text.append(clean_match)
        
        if extracted_text:
            # Join and deduplicate
            unique_text = []
            seen = set()
            for text in extracted_text:
                # Skip very similar text to avoid duplicates
                text_lower = text.lower()
                if not any(text_lower in seen_text for seen_text in seen):
                    unique_text.append(text)
                    seen.add(text_lower)
            
            result = ' '.join(unique_text)
            if len(result) > 100:
                logger.info(f"✅ Binary extraction found {len(result)} characters")
                return result
        
        logger.warning("⚠️ Binary extraction found no meaningful text")
        return ""
        
    except Exception as e:
        logger.warning(f"⚠️ Binary extraction failed: {e}")
        return ""

def extract_text_from_doc(content):
    """Extract text from DOC files using antiword."""
    try:
        # Create a temporary file for the DOC content
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as temp_file:
            temp_file.write(content)
            temp_file.flush()
            temp_doc_path = temp_file.name
        
        try:
            # Use antiword to extract text from DOC file
            result = subprocess.run([
                'antiword', temp_doc_path
            ], capture_output=True, text=True, timeout=30)
            
            if result.returncode == 0:
                return result.stdout.strip()
            else:
                logger.warning(f"⚠️ Antiword extraction failed: {result.stderr}")
                return ""
                
        finally:
            # Clean up the temporary DOC file
            if os.path.exists(temp_doc_path):
                os.unlink(temp_doc_path)
                
    except subprocess.TimeoutExpired:
        logger.warning("⚠️ Antiword extraction timed out")
        return ""
    except FileNotFoundError:
        logger.warning("⚠️ Antiword not found - please install with: brew install antiword")
        return ""
    except Exception as e:
        logger.warning(f"⚠️ Error with antiword extraction: {e}")
        return ""

def extract_text_from_file(content, mime_type):
    """Extract text from file content based on MIME type."""
    if 'wordprocessingml' in mime_type:  # DOCX
        return extract_text_from_docx(content)
    elif 'pdf' in mime_type:  # PDF
        return extract_text_from_pdf(content)
    elif 'presentation' in mime_type or 'powerpoint' in mime_type:  # PPT/PPTX
        return extract_text_from_ppt(content, mime_type)
    elif 'msword' in mime_type:  # DOC
        return extract_text_from_doc(content)
    else:
        logger.warning(f"⚠️ Unsupported MIME type: {mime_type}")
        return ""

def process_file(file_metadata, drive_service, openai_client, index):
    """Process a single file and create one chunk."""
    file_id = file_metadata['id']
    file_name = file_metadata['name']
    mime_type = file_metadata.get('mimeType', '')
    
    logger.info(f"🔄 Processing: {file_name}")
    logger.info(f"📄 File type: {mime_type}")
    
    try:
        # Download file
        content = download_file(drive_service, file_id)
        if not content:
            logger.error(f"❌ Failed to download {file_name}")
            return False
        
        logger.info(f"✅ Downloaded {len(content)} bytes")
        
        # Extract text
        text = extract_text_from_file(content, mime_type)
        
        if not text:
            logger.warning(f"⚠️ No text extracted from {file_name}")
            return False
        
        logger.info(f"✅ Extracted {len(text)} characters")
        
        # Take first 800 characters as one chunk
        chunk_text = text[:800] + "..." if len(text) > 800 else text
        logger.info(f"📄 Created chunk with {len(chunk_text)} characters")
        
        # Generate embedding
        logger.info("🧠 Generating embedding...")
        response = openai_client.embeddings.create(
            model="text-embedding-3-small",
            input=chunk_text
        )
        
        embedding = response.data[0].embedding
        logger.info(f"✅ Generated embedding with {len(embedding)} dimensions")
        
        # Prepare vector
        vector = {
            'id': f"{file_id}_single",
            'values': embedding,
            'metadata': {
                'text': chunk_text,
                'file_id': file_id,
                'file_name': file_name,
                'folder_path': 'site',
                'mime_type': mime_type,
                'modified_time': file_metadata.get('modifiedTime', '')
            }
        }
        
        # Upsert to Pinecone
        logger.info("📤 Upserting to Pinecone...")
        index.upsert(vectors=[vector], namespace='site')
        logger.info("✅ Successfully upserted to Pinecone")
        
        return True
        
    except Exception as e:
        logger.error(f"❌ Error processing {file_name}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return False

def test_query(index, openai_client):
    """Test query to verify ingestion worked."""
    logger.info("🔍 Testing query...")
    
    try:
        query = "safety training requirements"
        response = openai_client.embeddings.create(
            model="text-embedding-3-small",
            input=query
        )
        query_embedding = response.data[0].embedding
        
        results = index.query(
            vector=query_embedding,
            top_k=5,
            namespace='site',
            include_metadata=True
        )
        
        logger.info(f"🔍 Query returned {len(results.matches)} results:")
        for match in results.matches:
            logger.info(f"  - {match.metadata['file_name']} (Score: {match.score:.4f})")
        
        return True
        
    except Exception as e:
        logger.error(f"❌ Error testing query: {e}")
        return False

def main():
    """Main function."""
    logger.info("🚀 Starting comprehensive file ingestion")
    
    try:
        drive_service, openai_client, index = setup_clients()
        
        folder_id = os.getenv('GDRIVE_FOLDER_ID')
        files = get_all_supported_files(drive_service, folder_id)
        
        if not files:
            logger.error("❌ No supported files found")
            return
        
        successful_uploads = 0
        failed_uploads = 0
        
        for i, file_metadata in enumerate(files, 1):
            logger.info(f"\n{'='*60}")
            logger.info(f"📄 Processing file {i}/{len(files)}: {file_metadata['name']}")
            logger.info(f"{'='*60}")
            
            if process_file(file_metadata, drive_service, openai_client, index):
                successful_uploads += 1
                logger.info(f"✅ Successfully processed {file_metadata['name']}")
            else:
                failed_uploads += 1
                logger.error(f"❌ Failed to process {file_metadata['name']}")
        
        # Test query
        logger.info(f"\n{'='*60}")
        logger.info("🧪 Testing query functionality")
        logger.info(f"{'='*60}")
        test_query(index, openai_client)
        
        # Summary
        logger.info(f"\n{'='*60}")
        logger.info("📊 COMPREHENSIVE INGESTION SUMMARY")
        logger.info(f"{'='*60}")
        logger.info(f"📁 Files attempted: {len(files)}")
        logger.info(f"✅ Files successful: {successful_uploads}")
        logger.info(f"❌ Files failed: {failed_uploads}")
        logger.info(f"{'='*60}")
        
    except Exception as e:
        logger.error(f"❌ Pipeline failed: {e}")
        import traceback
        logger.error(traceback.format_exc())

if __name__ == "__main__":
    main()
