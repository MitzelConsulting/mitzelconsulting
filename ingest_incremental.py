#!/usr/bin/env python3
"""
Incremental ingestion script that only processes new or modified files.
Uses a manifest file to track processed files and their modification times.
"""

import os
import json
import logging
from datetime import datetime
from dotenv import load_dotenv
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from openai import OpenAI
from pinecone import Pinecone
from docx import Document
import io
import PyPDF2
from pptx import Presentation
import subprocess
import tempfile
import re

load_dotenv('.env.local')

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('ingest_incremental.log')
    ]
)
logger = logging.getLogger(__name__)

MANIFEST_FILE = 'ingestion_manifest.json'

def load_manifest():
    """Load the ingestion manifest file."""
    if os.path.exists(MANIFEST_FILE):
        try:
            with open(MANIFEST_FILE, 'r') as f:
                return json.load(f)
        except Exception as e:
            logger.warning(f"⚠️ Could not load manifest file: {e}")
    return {}

def save_manifest(manifest):
    """Save the ingestion manifest file."""
    try:
        with open(MANIFEST_FILE, 'w') as f:
            json.dump(manifest, f, indent=2)
        logger.info(f"💾 Manifest saved with {len(manifest)} entries")
    except Exception as e:
        logger.error(f"❌ Could not save manifest file: {e}")

def setup_clients():
    """Initialize clients."""
    logger.info("🔧 Setting up clients...")
    
    credentials = service_account.Credentials.from_service_account_file(
        'service-account.json',
        scopes=['https://www.googleapis.com/auth/drive.readonly']
    )
    drive_service = build('drive', 'v3', credentials=credentials)
    
    openai_client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
    
    pc = Pinecone(api_key=os.getenv('PINECONE_API_KEY'))
    index = pc.Index('mizelconsulting')
    
    logger.info("✅ Clients initialized")
    return drive_service, openai_client, index

def get_all_supported_files(drive_service, folder_id):
    """Get all supported files from the folder."""
    logger.info(f"📁 Getting all supported files from folder: {folder_id}")
    
    try:
        # Query for all supported file types
        supported_types = [
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # DOCX
            'application/pdf',  # PDF
            'application/vnd.ms-powerpoint',  # PPT
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # PPTX
            'application/msword'  # DOC
        ]
        
        # Build query for supported file types
        type_conditions = [f"mimeType='{mime_type}'" for mime_type in supported_types]
        type_query = " or ".join(type_conditions)
        full_query = f"'{folder_id}' in parents and trashed=false and ({type_query})"
        
        results = drive_service.files().list(
            q=full_query,
            fields="files(id, name, mimeType, modifiedTime, size)"
        ).execute()
        
        files = results.get('files', [])
        logger.info(f"📊 Found {len(files)} supported files")
        
        return files
        
    except Exception as e:
        logger.error(f"❌ Error getting files: {e}")
        return []

def get_new_or_modified_files(files, manifest):
    """Filter files to only include new or modified ones."""
    new_files = []
    modified_files = []
    
    for file in files:
        file_id = file['id']
        modified_time = file.get('modifiedTime', '')
        
        if file_id not in manifest:
            new_files.append(file)
            logger.info(f"🆕 New file: {file['name']}")
        elif manifest[file_id].get('modifiedTime') != modified_time:
            modified_files.append(file)
            logger.info(f"🔄 Modified file: {file['name']}")
    
    logger.info(f"📊 New files: {len(new_files)}")
    logger.info(f"📊 Modified files: {len(modified_files)}")
    
    return new_files + modified_files

def download_file(drive_service, file_id):
    """Download file content from Google Drive."""
    try:
        request = drive_service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while done is False:
            status, done = downloader.next_chunk()
        
        return fh.getvalue()
    except Exception as e:
        logger.error(f"❌ Error downloading file: {e}")
        return None

def extract_text_from_docx(content):
    """Extract text from DOCX content."""
    try:
        doc = Document(io.BytesIO(content))
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        return text.strip()
    except Exception as e:
        logger.error(f"❌ Error extracting DOCX text: {e}")
        return ""

def extract_text_from_pdf(content):
    """Extract text from PDF content."""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"❌ Error extracting PDF text: {e}")
        return ""

def extract_text_from_ppt(content, mime_type):
    """Extract text from PowerPoint files."""
    try:
        if 'presentationml' in mime_type:  # PPTX files
            prs = Presentation(io.BytesIO(content))
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += f"Slide {slide_num}: {shape.text.strip()}\n"
            return text.strip()
        else:  # PPT files (older format)
            # Try LibreOffice first, then fall back to binary extraction
            text = extract_text_from_old_ppt(content)
            if not text:
                text = extract_text_from_ppt_binary(content)
            return text
        
    except Exception as e:
        logger.error(f"❌ Error extracting PPT text: {e}")
        return ""

def extract_text_from_old_ppt(content):
    """Extract text from old PPT files using LibreOffice command line."""
    try:
        # Create a temporary file for the PPT content
        with tempfile.NamedTemporaryFile(suffix='.ppt', delete=False) as temp_file:
            temp_file.write(content)
            temp_file.flush()
            temp_ppt_path = temp_file.name
        
        try:
            # Use LibreOffice to convert PPT to text
            result = subprocess.run([
                '/Applications/LibreOffice.app/Contents/MacOS/soffice', '--headless', '--convert-to', 'txt', 
                temp_ppt_path
            ], capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                # LibreOffice creates the converted file in the current working directory
                temp_filename = os.path.basename(temp_ppt_path)
                txt_filename = temp_filename.replace('.ppt', '.txt')
                txt_path = os.path.join(os.getcwd(), txt_filename)
                
                if os.path.exists(txt_path):
                    with open(txt_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                    os.unlink(txt_path)
                    return text.strip()
                else:
                    return ""
            else:
                return ""
                
        finally:
            if os.path.exists(temp_ppt_path):
                os.unlink(temp_ppt_path)
                
    except Exception as e:
        return ""

def extract_text_from_ppt_binary(content):
    """Extract readable text directly from PPT binary content."""
    try:
        # Convert to string and look for readable text patterns
        text_content = content.decode('latin-1', errors='ignore')
        
        # Look for common text patterns in PPT files
        patterns = [
            r'[\x20-\x7E]{15,}',  # Printable ASCII characters, at least 15 chars
            r'[A-Za-z\s]{25,}',   # Letters and spaces, at least 25 chars
        ]
        
        extracted_text = []
        for pattern in patterns:
            matches = re.findall(pattern, text_content)
            for match in matches:
                # Clean up the match
                clean_match = match.strip()
                # Filter out binary noise and keep meaningful text
                if (len(clean_match) > 15 and 
                    not any(char in clean_match for char in ['\x00', '\x01', '\x02', '\x03']) and
                    any(c.isalpha() for c in clean_match) and  # Must contain letters
                    len([c for c in clean_match if c.isalpha()]) > 5):  # At least 5 letters
                    extracted_text.append(clean_match)
        
        if extracted_text:
            # Join and deduplicate
            unique_text = []
            seen = set()
            for text in extracted_text:
                # Skip very similar text to avoid duplicates
                text_lower = text.lower()
                if not any(text_lower in seen_text for seen_text in seen):
                    unique_text.append(text)
                    seen.add(text_lower)
            
            result = ' '.join(unique_text)
            if len(result) > 100:
                return result
        
        return ""
        
    except Exception as e:
        return ""

def extract_text_from_doc(content):
    """Extract text from DOC files using antiword."""
    try:
        # Create a temporary file for the DOC content
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as temp_file:
            temp_file.write(content)
            temp_file.flush()
            temp_doc_path = temp_file.name
        
        try:
            # Use antiword to extract text from DOC file
            result = subprocess.run([
                'antiword', temp_doc_path
            ], capture_output=True, text=True, timeout=30)
            
            if result.returncode == 0:
                return result.stdout.strip()
            else:
                return ""
                
        finally:
            if os.path.exists(temp_doc_path):
                os.unlink(temp_doc_path)
                
    except Exception as e:
        return ""

def extract_text_from_file(content, mime_type):
    """Extract text from file content based on MIME type."""
    if 'wordprocessingml' in mime_type:  # DOCX
        return extract_text_from_docx(content)
    elif 'pdf' in mime_type:  # PDF
        return extract_text_from_pdf(content)
    elif 'presentation' in mime_type or 'powerpoint' in mime_type:  # PPT/PPTX
        return extract_text_from_ppt(content, mime_type)
    elif 'msword' in mime_type:  # DOC
        return extract_text_from_doc(content)
    else:
        logger.warning(f"⚠️ Unsupported MIME type: {mime_type}")
        return ""

def process_file(file_metadata, drive_service, openai_client, index):
    """Process a single file and create one chunk."""
    file_id = file_metadata['id']
    file_name = file_metadata['name']
    mime_type = file_metadata.get('mimeType', '')
    
    logger.info(f"🔄 Processing: {file_name}")
    
    try:
        # Download file
        content = download_file(drive_service, file_id)
        if not content:
            logger.error(f"❌ Failed to download {file_name}")
            return False
        
        logger.info(f"✅ Downloaded {len(content)} bytes")
        
        # Extract text
        text = extract_text_from_file(content, mime_type)
        
        if not text:
            logger.warning(f"⚠️ No text extracted from {file_name}")
            return False
        
        logger.info(f"✅ Extracted {len(text)} characters")
        
        # Take first 800 characters as one chunk
        chunk_text = text[:800] + "..." if len(text) > 800 else text
        logger.info(f"📄 Created chunk with {len(chunk_text)} characters")
        
        # Generate embedding
        logger.info("🧠 Generating embedding...")
        response = openai_client.embeddings.create(
            model="text-embedding-3-small",
            input=chunk_text
        )
        
        embedding = response.data[0].embedding
        logger.info(f"✅ Generated embedding with {len(embedding)} dimensions")
        
        # Prepare vector
        vector = {
            'id': f"{file_id}_single",
            'values': embedding,
            'metadata': {
                'text': chunk_text,
                'file_id': file_id,
                'file_name': file_name,
                'folder_path': 'site',
                'mime_type': mime_type,
                'modified_time': file_metadata.get('modifiedTime', '')
            }
        }
        
        # Upsert to Pinecone
        logger.info("📤 Upserting to Pinecone...")
        index.upsert(vectors=[vector], namespace='site')
        logger.info("✅ Successfully upserted to Pinecone")
        
        return True
        
    except Exception as e:
        logger.error(f"❌ Error processing {file_name}: {e}")
        return False

def main():
    """Main function."""
    logger.info("🚀 Starting incremental file ingestion")
    
    try:
        # Load manifest
        manifest = load_manifest()
        logger.info(f"📋 Loaded manifest with {len(manifest)} previously processed files")
        
        drive_service, openai_client, index = setup_clients()
        
        folder_id = os.getenv('GDRIVE_FOLDER_ID')
        all_files = get_all_supported_files(drive_service, folder_id)
        
        if not all_files:
            logger.error("❌ No supported files found")
            return
        
        # Get only new or modified files
        files_to_process = get_new_or_modified_files(all_files, manifest)
        
        if not files_to_process:
            logger.info("✅ No new or modified files to process")
            return
        
        successful_uploads = 0
        failed_uploads = 0
        
        for i, file_metadata in enumerate(files_to_process, 1):
            logger.info(f"\n{'='*60}")
            logger.info(f"📄 Processing file {i}/{len(files_to_process)}: {file_metadata['name']}")
            logger.info(f"{'='*60}")
            
            if process_file(file_metadata, drive_service, openai_client, index):
                successful_uploads += 1
                logger.info(f"✅ Successfully processed {file_metadata['name']}")
                
                # Update manifest
                manifest[file_metadata['id']] = {
                    'name': file_metadata['name'],
                    'modifiedTime': file_metadata.get('modifiedTime', ''),
                    'processedAt': datetime.now().isoformat()
                }
            else:
                failed_uploads += 1
                logger.error(f"❌ Failed to process {file_metadata['name']}")
        
        # Save updated manifest
        save_manifest(manifest)
        
        # Summary
        logger.info(f"\n{'='*60}")
        logger.info("📊 INCREMENTAL INGESTION SUMMARY")
        logger.info(f"{'='*60}")
        logger.info(f"📁 Files attempted: {len(files_to_process)}")
        logger.info(f"✅ Files successful: {successful_uploads}")
        logger.info(f"❌ Files failed: {failed_uploads}")
        logger.info(f"📈 Success rate: {(successful_uploads/len(files_to_process)*100):.1f}%")
        logger.info(f"{'='*60}")
        
    except Exception as e:
        logger.error(f"❌ Pipeline failed: {e}")
        import traceback
        logger.error(traceback.format_exc())

if __name__ == "__main__":
    main()

